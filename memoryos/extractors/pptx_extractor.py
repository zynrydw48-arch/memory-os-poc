from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from memoryos.extractors.result import ExtractionResult

MAX_EMBEDDED_IMAGES = 5  # bounds processing time on image-heavy decks


def extract_pptx(path) -> ExtractionResult:
    presentation = Presentation(path)
    text_parts = []
    embedded_images: list[bytes] = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                text_parts.append(shape.text_frame.text.strip())
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text.strip())

            if (
                shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                and len(embedded_images) < MAX_EMBEDDED_IMAGES
            ):
                embedded_images.append(shape.image.blob)

    return ExtractionResult(
        text="\n".join(text_parts),
        embedded_images=embedded_images,
        structural_metadata={"slide_count": len(presentation.slides)},
    )
